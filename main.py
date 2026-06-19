import io
import os
import copy
import openpyxl
import streamlit as st
from collections import defaultdict
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Default PPT template bundled with the app. Shape mutations target shapes by their
# exact names in this file (see CLAUDE.md → "Shape targeting"), so it is the
# reference layout. Users only upload the Excel; the template loads automatically.
TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'PPTmuestra.pptx')

# Template geometry (EMU) — measured from PPTmuestra.pptx
CIRCLE_X = 11337402
CIRCLE_Y0 = 1203044   # top of first circle
CIRCLE_STEP = 731520  # vertical step per row
CIRCLE_SZ = 288000

FILL_RED  = f'<a:solidFill xmlns:a="{A_NS}"><a:srgbClr val="FF0000"/></a:solidFill>'
FILL_GRAY = f'<a:solidFill xmlns:a="{A_NS}"><a:schemeClr val="bg1"><a:lumMod val="65000"/></a:schemeClr></a:solidFill>'

# Separator rectangles (row indicators)
RECT_X    = 5313211
RECT_Y0   = 1614107   # top of first rect (row 1)
RECT_STEP = 716400    # = 1.99cm, matches data row height
RECT_SZ_W = 135089
RECT_SZ_H = 137865

FILL_TRANSVERSAL = f'<a:solidFill xmlns:a="{A_NS}"><a:srgbClr val="00B491"/></a:solidFill>'
FILL_ESPECIFICO  = f'<a:solidFill xmlns:a="{A_NS}"><a:schemeClr val="tx1"><a:lumMod val="75000"/><a:lumOff val="25000"/></a:schemeClr></a:solidFill>'

# Career sidebar (Tabla 8) — section cards on top, careers of the active section
# in the bottom slot. The active section/career are emphasised. Row layout is read
# from the template at runtime (see update_career_list); these are the row roles.
SIDEBAR_LIST_ROW  = 4            # last row of Tabla 8 holds the active section's career list
SIDEBAR_MENU_ROWS = 4            # rows 0..3 are the section cards
SIDEBAR_PARA_H    = 265045       # EMU per career line in the list row (font metric)
ACTIVE_COLOR      = RGBColor(0x00, 0x00, 0x00)  # black; active item stands out via bold


# ── Excel reading ────────────────────────────────────────────────────────────

def read_excel(excel_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))

    ws = wb['datos']
    headers = [c.value for c in ws[1]]
    datos = [
        dict(zip(headers, [c.value for c in row]))
        for row in ws.iter_rows(min_row=2)
        if any(c.value for c in row)
    ]

    ws = wb['metadata']
    headers = [c.value for c in ws[1]]
    meta = {}
    for row in ws.iter_rows(min_row=2):
        if any(c.value for c in row):
            d = dict(zip(headers, [c.value for c in row]))
            meta[(d['carrera'], d['segmento'], d['modalidad'])] = d

    secciones = read_secciones(wb)

    return datos, meta, secciones


def read_secciones(wb):
    """Read the `secciones` sheet (columns: seccion, carrera) — the single source
    of truth for which careers belong to which section, and in what order.

    Returns a dict with:
      - order:     section names, in first-seen order (the sidebar menu).
      - careers:   {section: [careers in sheet order]}.
      - by_career: {career: section}, for reverse lookup.
    """
    order = []
    careers = {}
    by_career = {}

    if 'secciones' not in wb.sheetnames:
        return {'order': order, 'careers': careers, 'by_career': by_career}

    ws = wb['secciones']
    headers = [c.value for c in ws[1]]
    for row in ws.iter_rows(min_row=2):
        if not any(c.value for c in row):
            continue
        d = dict(zip(headers, [c.value for c in row]))
        seccion = (d.get('seccion') or '').strip() if isinstance(d.get('seccion'), str) else d.get('seccion')
        carrera = (d.get('carrera') or '').strip() if isinstance(d.get('carrera'), str) else d.get('carrera')
        if not seccion or not carrera:
            continue
        if seccion not in careers:
            careers[seccion] = []
            order.append(seccion)
        if carrera not in careers[seccion]:
            careers[seccion].append(carrera)
        by_career[carrera] = seccion

    return {'order': order, 'careers': careers, 'by_career': by_career}


# ── Slide cloning ────────────────────────────────────────────────────────────

def clone_slide(output_prs, template_slide):
    """Append a clone of template_slide to output_prs, skipping OLE objects."""
    new_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
    sp_tree = new_slide.shapes._spTree

    # Keep first 2 structural children (nvGrpSpPr, grpSpPr), remove the rest
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)

    for shape in template_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            continue
        sp_tree.append(copy.deepcopy(shape._element))

    return new_slide


# ── Shape helpers ────────────────────────────────────────────────────────────

def find_shape(slide, name):
    return next((s for s in slide.shapes if s.name == name), None)


def set_cell(cell, text, bold=None, color=None):
    """Overwrite cell text preserving first-run formatting.

    Optionally force `bold` and/or font `color` (an RGBColor) on the run —
    used to emphasise the active item in the career sidebar.
    """
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    if para.runs:
        run = para.runs[0]
        run.text = str(text) if text is not None else ''
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        run = para.add_run()
        run.text = str(text) if text is not None else ''
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def set_cell_lines(cell, lines, active=None, active_color=None):
    """Fill a cell with one paragraph per line, cloning the first paragraph so
    every line keeps the template's run formatting.

    If `active` matches a line, that line is bolded (and recoloured when
    `active_color` is given). Used for the career list in the sidebar.
    """
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    lines = [str(l) for l in (lines or [''])]

    # Match paragraph count to line count by cloning the first paragraph.
    while len(tf.paragraphs) < len(lines):
        clone = copy.deepcopy(tf.paragraphs[0]._p)
        tf.paragraphs[-1]._p.addnext(clone)
    while len(tf.paragraphs) > len(lines):
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)

    for para, line in zip(tf.paragraphs, lines):
        if para.runs:
            run = para.runs[0]
            run.text = line
            for r in para.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            run = para.add_run()
            run.text = line
        is_active = active is not None and line == active
        run.font.bold = True if is_active else None
        if is_active and active_color is not None:
            run.font.color.rgb = active_color


# ── Per-slide updaters ───────────────────────────────────────────────────────

def update_header(slide, carrera, segmento, modalidad, des_carrera, des_prom, pg, total):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top / 360000 < 5:
            if 'Título' in shape.name or 'Titulo' in shape.name:
                text = (
                    f"{carrera} – {segmento} – {modalidad}"
                    f"  |  Deserción Carrera: {des_carrera}"
                    f"   |   Deserción Prom.: {des_prom}\t    \t      ({pg} de {total})"
                )
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = text
                break


def update_comments(slide, comentarios):
    """comentarios can be a string with \n separating paragraphs."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top / 360000 > 10:
            if 'Título' in shape.name or 'Titulo' in shape.name:
                tf = shape.text_frame
                lines = str(comentarios).split('\n') if comentarios else ['']

                # Ensure enough paragraphs
                while len(tf.paragraphs) < len(lines):
                    # Clone last paragraph
                    last_p = tf.paragraphs[-1]._p
                    new_p = copy.deepcopy(last_p)
                    last_p.addnext(new_p)

                # Remove excess paragraphs
                while len(tf.paragraphs) > len(lines):
                    tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)

                for para, line in zip(tf.paragraphs, lines):
                    if para.runs:
                        para.runs[0].text = line
                        for r in para.runs[1:]:
                            r._r.getparent().remove(r._r)
                    else:
                        para.add_run().text = line
                break


def update_main_table(slide, rows):
    """rows: list of (pain, hallazgo, accion) tuples."""
    shape = find_shape(slide, 'Tabla 87')
    if not shape or not shape.has_table:
        return
    table = shape.table

    # Adjust row count (header is row 0, data starts at row 1)
    while len(table.rows) - 1 < len(rows):
        new_tr = copy.deepcopy(table.rows[1]._tr)  # clone first data row
        table._tbl.append(new_tr)

    while len(table.rows) - 1 > len(rows):
        table._tbl.remove(table.rows[len(table.rows) - 1]._tr)

    for i, (pain, hallazgo, accion) in enumerate(rows):
        r = table.rows[i + 1]
        set_cell(r.cells[0], pain)
        set_cell(r.cells[1], hallazgo)
        set_cell(r.cells[2], accion)


def update_circles(slide, niveles):
    """niveles: list of 'profundizar' | 'no_accionable'."""
    circles = [s for s in slide.shapes if 'Elipse' in s.name]
    if not circles:
        return

    template_sp = circles[0]._element

    # Remove excess
    for c in circles[len(niveles):]:
        c._element.getparent().remove(c._element)

    # Add missing (clone first circle)
    circles = [s for s in slide.shapes if 'Elipse' in s.name]
    while len(circles) < len(niveles):
        new_sp = copy.deepcopy(template_sp)
        slide.shapes._spTree.append(new_sp)
        circles = [s for s in slide.shapes if 'Elipse' in s.name]

    # Position + color
    for i, circle in enumerate(circles):
        spPr = circle._element.find(f'{{{P_NS}}}spPr')
        if spPr is None:
            continue
        xfrm = spPr.find(f'{{{A_NS}}}xfrm')
        if xfrm is not None:
            off = xfrm.find(f'{{{A_NS}}}off')
            if off is not None:
                off.set('x', str(CIRCLE_X))
                off.set('y', str(CIRCLE_Y0 + i * CIRCLE_STEP))

        # Replace fill
        for tag in [f'{{{A_NS}}}solidFill', f'{{{A_NS}}}noFill', f'{{{A_NS}}}gradFill']:
            for el in spPr.findall(tag):
                spPr.remove(el)

        nivel = niveles[i] if i < len(niveles) else 'no_accionable'
        fill_xml = FILL_RED if (nivel or '').lower() == 'profundizar' else FILL_GRAY
        new_fill = etree.fromstring(fill_xml)
        geom = spPr.find(f'{{{A_NS}}}prstGeom')
        if geom is not None:
            geom.addnext(new_fill)
        else:
            spPr.append(new_fill)


def update_rects(slide, separadores):
    """separadores: list of 'transversal' | 'especifico'"""
    rects = [s for s in slide.shapes if 'Rect' in s.name and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    if not rects:
        return

    template_sp = rects[0]._element

    # Remove excess
    for r in rects[len(separadores):]:
        r._element.getparent().remove(r._element)

    # Add missing
    rects = [s for s in slide.shapes if 'Rect' in s.name and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    while len(rects) < len(separadores):
        new_sp = copy.deepcopy(template_sp)
        slide.shapes._spTree.append(new_sp)
        rects = [s for s in slide.shapes if 'Rect' in s.name and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]

    # Position + color
    for i, rect in enumerate(rects):
        spPr = rect._element.find(f'{{{P_NS}}}spPr')
        if spPr is None:
            continue
        xfrm = spPr.find(f'{{{A_NS}}}xfrm')
        if xfrm is not None:
            off = xfrm.find(f'{{{A_NS}}}off')
            ext = xfrm.find(f'{{{A_NS}}}ext')
            if off is not None:
                off.set('x', str(RECT_X))
                off.set('y', str(RECT_Y0 + i * RECT_STEP))
            if ext is not None:
                ext.set('cx', str(RECT_SZ_W))
                ext.set('cy', str(RECT_SZ_H))

        for tag in [f'{{{A_NS}}}solidFill', f'{{{A_NS}}}noFill', f'{{{A_NS}}}gradFill']:
            for el in spPr.findall(tag):
                spPr.remove(el)

        sep = (separadores[i] or '').lower() if i < len(separadores) else 'transversal'
        fill_xml = FILL_ESPECIFICO if sep == 'especifico' else FILL_TRANSVERSAL
        new_fill = etree.fromstring(fill_xml)
        geom = spPr.find(f'{{{A_NS}}}prstGeom')
        if geom is not None:
            geom.addnext(new_fill)
        else:
            spPr.append(new_fill)


def update_career_list(slide, carrera, seccion, secciones):
    """Render the career sidebar (Tabla 8) for the active career.

    Layout — "inline expansion": the careers of the active section appear directly
    below that section's card, pushing the sections below it down. Example (active
    section = Salud):

        Carreras Masivas
        Carreras Faco
        Salud              ← active section card (bolded + recoloured)
          Medicina         ← active career (bolded + arrow)
          Enfermería
          Obstetricia
        Carreras No Masivas

    Implementation: the template's single list row (SIDEBAR_LIST_ROW) is filled with
    the active section's careers and its `<a:tr>` is moved to sit right after the
    active section's card. The highlight bar + arrow (Grupo 20) point at the active
    career. Geometry (row heights, group coordinate offset) is read at runtime, so
    the layout survives template edits.
    """
    shape = find_shape(slide, 'Tabla 8')
    if not shape or not shape.has_table:
        return
    table = shape.table
    if len(table.rows) <= SIDEBAR_LIST_ROW:
        return

    order = secciones.get('order', [])
    card_heights = [table.rows[j].height for j in range(SIDEBAR_MENU_ROWS)]

    # 1) Section cards. Extra rows are blanked. Remember the active card's index.
    active_k = None
    for j in range(SIDEBAR_MENU_ROWS):
        cell = table.rows[j].cells[0]
        if j < len(order):
            sec = order[j]
            is_active = sec == seccion
            if is_active:
                active_k = j
            set_cell(cell, sec, bold=is_active,
                     color=ACTIVE_COLOR if is_active else None)
        else:
            set_cell(cell, '', bold=False)

    # 2) Fill the list row with the active section's careers.
    careers = secciones.get('careers', {}).get(seccion, [])
    set_cell_lines(table.rows[SIDEBAR_LIST_ROW].cells[0], careers,
                   active=carrera, active_color=ACTIVE_COLOR)

    # 3) Move the list row inline — right after the active section's card.
    trs = table._tbl.findall(f'{{{A_NS}}}tr')
    list_tr = trs[SIDEBAR_LIST_ROW]
    if active_k is not None:
        trs[active_k].addnext(list_tr)
        rows_before = active_k + 1          # cards 0..active_k precede the list
    else:
        rows_before = SIDEBAR_MENU_ROWS     # unknown section → leave list at bottom

    # 4) Point the highlight bar + arrow (inside Grupo 20) at the active career.
    list_top = shape.top + sum(card_heights[:rows_before])
    try:
        idx = careers.index(carrera)
    except ValueError:
        idx = 0
    center_y = int(list_top + (idx + 0.5) * SIDEBAR_PARA_H)

    grupo = find_shape(slide, 'Grupo 20')
    if grupo is None:
        return

    # Convert an absolute Y to the group's child coordinate space: a child renders
    # at off + (child_local - chOff), so child_local = abs + (chOff - off).
    group_xfrm = grupo._element.find(f'{{{P_NS}}}grpSpPr/{{{A_NS}}}xfrm')
    local_offset = 0
    if group_xfrm is not None:
        g_off = group_xfrm.find(f'{{{A_NS}}}off')
        g_choff = group_xfrm.find(f'{{{A_NS}}}chOff')
        if g_off is not None and g_choff is not None:
            local_offset = int(g_choff.get('y')) - int(g_off.get('y'))

    for child in grupo.shapes:
        is_rect  = 'Rect' in child.name
        is_arrow = 'Flecha' in child.name or 'flecha' in child.name
        if not (is_rect or is_arrow):
            continue

        new_local_y = center_y - child.height // 2 + local_offset

        spPr = child._element.find(f'{{{P_NS}}}spPr')
        if spPr is None:
            continue
        xfrm = spPr.find(f'{{{A_NS}}}xfrm')
        if xfrm is None:
            continue
        off = xfrm.find(f'{{{A_NS}}}off')
        if off is not None:
            off.set('y', str(new_local_y))


def update_stats(slide, muestra, desertores, alumnos):
    shape = find_shape(slide, 'Tabla 45')
    if not shape or not shape.has_table:
        return
    t = shape.table
    set_cell(t.rows[0].cells[1], muestra)
    set_cell(t.rows[1].cells[1], desertores)
    set_cell(t.rows[2].cells[1], alumnos)


def update_var_table(slide, var_des_prom, var_periodo_ant):
    shape = find_shape(slide, 'Table 11')
    if not shape or not shape.has_table:
        return
    t = shape.table
    set_cell(t.rows[1].cells[0], var_des_prom)
    set_cell(t.rows[1].cells[1], var_periodo_ant)


# ── Main generator ───────────────────────────────────────────────────────────

def generate_ppt(template_bytes, datos, metadata, secciones):
    template_prs = Presentation(io.BytesIO(template_bytes))
    template_slide = template_prs.slides[0]

    output_prs = Presentation(io.BytesIO(template_bytes))

    groups = defaultdict(list)
    for row in datos:
        key = (row.get('carrera'), row.get('segmento'), row.get('modalidad'))
        groups[key].append(row)

    for i, (key, rows) in enumerate(groups.items()):
        carrera, segmento, modalidad = key
        meta = metadata.get(key, {})
        slide = output_prs.slides[0] if i == 0 else clone_slide(output_prs, template_slide)

        update_header(slide, carrera, segmento, modalidad,
                      meta.get('desercion_carrera', ''),
                      meta.get('desercion_prom', ''),
                      1, 1)

        table_rows = [(r.get('pain'), r.get('hallazgo'), r.get('accion')) for r in rows]
        update_main_table(slide, table_rows)

        niveles     = [r.get('nivel', 'no_accionable') for r in rows]
        separadores = [r.get('separador', 'transversal') for r in rows]
        update_circles(slide, niveles)
        update_rects(slide, separadores)

        # Section comes from the `secciones` sheet (source of truth); fall back
        # to the datos `tipo_carrera` column only if the career is unmapped.
        seccion = secciones['by_career'].get(carrera) or rows[0].get('tipo_carrera', '')
        update_career_list(slide, carrera, seccion, secciones)

        update_stats(slide, meta.get('muestra'), meta.get('desertores'), meta.get('alumnos'))
        update_var_table(slide, meta.get('var_des_prom'), meta.get('var_periodo_ant'))
        update_comments(slide, meta.get('comentarios'))

    buf = io.BytesIO()
    output_prs.save(buf)
    return buf.getvalue()


# ── Streamlit UI ─────────────────────────────────────────────────────────────

def load_template_bytes(uploaded):
    """Return the PPT template bytes.

    Uses the uploaded file when the user provides one; otherwise falls back to
    the bundled PPTmuestra.pptx so the user only has to upload the Excel.
    Returns None if no template is available.
    """
    if uploaded is not None:
        return uploaded.read()
    if os.path.exists(TEMPLATE_PATH):
        with open(TEMPLATE_PATH, 'rb') as f:
            return f.read()
    return None


_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

:root {
  --brand: #00B491;
  --brand-deep: #00897B;
  --brand-deeper: #006F63;
  --ink: #0F2D28;
  --muted: #3A4D49;
  --surface: #FFFFFF;
  --tint: #E4F4F0;
  --line: #DCE6E3;
}

html, body, .stApp, [class*="css"] {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
}

/* Clean chrome for a deployed, non-technical audience */
#MainMenu, footer { visibility: hidden; }
[data-testid="stHeader"] { background: transparent; }

.block-container {
  max-width: 820px;
  padding-top: 2.5rem;
  padding-bottom: 4rem;
}

/* Hero */
.app-hero { margin-bottom: 1.75rem; }
.app-pill {
  display: inline-block;
  font-size: 0.78rem;
  font-weight: 600;
  color: var(--brand-deep);
  background: var(--tint);
  padding: 0.3rem 0.7rem;
  border-radius: 999px;
}
.app-title {
  margin: 0.75rem 0 0.4rem;
  font-size: clamp(2.1rem, 5vw, 3rem);
  font-weight: 700;
  letter-spacing: -0.025em;
  line-height: 1.05;
  color: var(--ink);
  text-wrap: balance;
}
.app-title::after {
  content: "";
  display: block;
  width: 56px;
  height: 4px;
  margin-top: 0.7rem;
  border-radius: 2px;
  background: var(--brand);
}
.app-sub {
  margin: 0;
  max-width: 60ch;
  font-size: 1.02rem;
  line-height: 1.55;
  color: var(--muted);
}

/* File uploader dropzone */
[data-testid="stFileUploaderDropzone"] {
  background: var(--surface);
  border: 1.5px dashed var(--line);
  border-radius: 14px;
  transition: border-color .2s ease, background .2s ease;
}
[data-testid="stFileUploaderDropzone"]:hover {
  border-color: var(--brand);
  background: #FBFEFD;
}

/* Buttons */
.stButton > button, .stDownloadButton > button {
  width: 100%;
  border-radius: 10px;
  font-weight: 600;
  padding: 0.6rem 1rem;
  border: none;
  transition: transform .12s ease, box-shadow .2s ease, background .2s ease;
}
.stButton > button[kind="primary"], .stDownloadButton > button[kind="primary"] {
  background: var(--brand-deep);
  color: #fff;
  box-shadow: 0 6px 16px -6px rgba(0, 137, 123, 0.55);
}
.stButton > button[kind="primary"]:hover, .stDownloadButton > button[kind="primary"]:hover {
  background: var(--brand-deeper);
  transform: translateY(-1px);
}
.stButton > button[kind="primary"]:active, .stDownloadButton > button[kind="primary"]:active {
  transform: translateY(0);
}

/* Expander */
[data-testid="stExpander"] {
  border: 1px solid var(--line);
  border-radius: 12px;
  background: var(--surface);
}

@media (prefers-reduced-motion: reduce) {
  .stButton > button, .stDownloadButton > button,
  [data-testid="stFileUploaderDropzone"] { transition: none; }
  .stButton > button[kind="primary"]:hover,
  .stDownloadButton > button[kind="primary"]:hover { transform: none; }
}
</style>
"""

st.set_page_config(page_title="AutoPPT", page_icon="📊", layout="centered")
st.markdown(_CSS, unsafe_allow_html=True)
st.markdown(
    """
    <div class="app-hero">
      <span class="app-pill">Generador de presentaciones</span>
      <h1 class="app-title">AutoPPT</h1>
      <p class="app-sub">Convertí tu Excel de hallazgos en una presentación lista para
      revisar. Subí el archivo y descargá el PPT generado.</p>
    </div>
    """,
    unsafe_allow_html=True,
)

with st.expander("📋 Formato del Excel requerido"):
    st.markdown("""
**Hoja `datos`** — una fila por pain/hallazgo:

| tipo_carrera | carrera | segmento | modalidad | pain | hallazgo | accion | nivel |
|---|---|---|---|---|---|---|---|
| Masivas | Ing. Electrónica | Nuevos | Presencial | Pagos... | Se le generó... | No aplica | no_accionable |

`nivel`: `profundizar` → círculo rojo · `no_accionable` → círculo gris

---

**Hoja `metadata`** — una fila por combinación carrera × segmento × modalidad:

| carrera | segmento | modalidad | desercion_carrera | desercion_prom | muestra | desertores | alumnos | comentarios | var_des_prom | var_periodo_ant |
|---|---|---|---|---|---|---|---|---|---|---|

`comentarios`: usá Alt+Enter dentro de la celda para múltiples líneas.

---

**Hoja `secciones`** — define a qué sección pertenece cada carrera (una fila por carrera):

| seccion | carrera |
|---|---|
| Carreras Masivas | Contabilidad |
| Salud | Medicina |

En cada slide, la barra lateral muestra las tarjetas de sección y, debajo, **solo las
carreras de la sección de esa carrera**, resaltando la actual. El orden de las filas
define el orden en que aparecen. Soporta hasta 5 secciones en el menú.
    """)

col1, col2 = st.columns(2)
with col1:
    excel_file = st.file_uploader("📄 Excel con datos", type=["xlsx"])
with col2:
    ppt_file = st.file_uploader("📑 Plantilla PPT (opcional)", type=["pptx"])

if ppt_file is not None:
    st.caption("📑 Usando la plantilla que subiste.")
elif os.path.exists(TEMPLATE_PATH):
    st.caption("📑 Usando la plantilla incluida: **PPTmuestra.pptx**.")
else:
    st.warning("No se encontró la plantilla incluida `PPTmuestra.pptx`. Subí una plantilla PPT.")

if excel_file:
    try:
        datos, metadata, secciones = read_excel(excel_file.read())
        keys = list({(r.get('carrera'), r.get('segmento'), r.get('modalidad')) for r in datos})
        st.info(f"Se generarán **{len(keys)} slides** para {len(datos)} filas de datos.")

        # Validación de la hoja `secciones`.
        if not secciones['order']:
            st.warning(
                "No se encontró la hoja **`secciones`** (o está vacía). "
                "La barra lateral de carreras no se podrá agrupar por sección."
            )
        else:
            sin_seccion = sorted({
                r.get('carrera') for r in datos
                if r.get('carrera') and r.get('carrera') not in secciones['by_career']
            })
            if sin_seccion:
                st.warning(
                    "Estas carreras de `datos` no están en la hoja `secciones` y "
                    "no se agruparán: " + ", ".join(sin_seccion)
                )
            if len(secciones['order']) > SIDEBAR_MENU_ROWS:
                st.warning(
                    f"Hay {len(secciones['order'])} secciones pero el template solo "
                    f"muestra {SIDEBAR_MENU_ROWS} en el menú. Las demás no aparecerán como tarjeta."
                )

        template_bytes = load_template_bytes(ppt_file)
        if template_bytes is None:
            st.error("No hay plantilla PPT disponible. Subí una para generar el PPT.")
        elif st.button("🚀 Generar PPT", type="primary"):
            with st.spinner("Generando..."):
                result = generate_ppt(template_bytes, datos, metadata, secciones)
            st.success("¡Listo!")
            st.download_button(
                label="⬇️ Descargar PPT generado",
                data=result,
                file_name="resultado.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
            )
    except Exception as e:
        st.error(f"Error: {e}")
        import traceback
        st.code(traceback.format_exc())
