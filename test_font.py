"""Check that every run of the generated deck uses CONTENT_FONT.

Run: python test_font.py
"""
import io

from pptx import Presentation
from pptx.oxml.ns import qn

import main


def test_font():
    datos, metadata, secciones = main.read_excel(open('ejemplo.xlsx', 'rb').read())
    out = main.generate_ppt(open(main.TEMPLATE_PATH, 'rb').read(), datos, metadata, secciones)
    prs = Presentation(io.BytesIO(out))

    faces = {l.get('typeface') for s in prs.slides for l in s._element.iter(qn('a:latin'))}
    assert faces == {main.CONTENT_FONT}, faces
    assert not [r for s in prs.slides for r in s._element.iter(qn('a:r'))
                if r.find(qn('a:rPr')) is None], 'runs without rPr inherit the theme font'


if __name__ == '__main__':
    test_font()
    print('ok')
