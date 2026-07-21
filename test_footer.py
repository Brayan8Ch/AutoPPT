"""Check the footer variation table: cell background per value, white text.

Run: python test_footer.py
"""
import io

from pptx import Presentation

import main


def test_var_table():
    datos, metadata, secciones = main.read_excel(open('ejemplo.xlsx', 'rb').read())

    key = next(k for k in metadata if k in {
        (r.get('carrera'), r.get('segmento'), r.get('modalidad')) for r in datos
    })
    metadata[key]['var_des_prom'] = 'Muy Encima'
    metadata[key]['var_periodo_ant'] = 'Debajo'

    out = main.generate_ppt(open(main.TEMPLATE_PATH, 'rb').read(), datos, metadata, secciones)
    prs = Presentation(io.BytesIO(out))

    found = False
    for slide in prs.slides:
        shape = main.find_shape(slide, 'Table 11')
        cells = shape.table.rows[1].cells
        if cells[0].text != 'Muy Encima':
            continue
        found = True
        assert cells[0].fill.fore_color.rgb == main.VAR_COLORS['muy encima']
        assert cells[1].fill.fore_color.rgb == main.VAR_COLORS['debajo']
        for cell in cells:
            assert cell.text_frame.paragraphs[0].runs[0].font.color.rgb == main.VAR_TEXT_COLOR
    assert found, 'no slide got the injected variation values'


def test_footer_matrix():
    prs = Presentation(main.TEMPLATE_PATH)
    slide = prs.slides[0]

    matched = main.update_footer_matrix(slide, 'Año 1', 'CGT 50/50')  # → Semipresencial
    assert matched
    t = main.find_shape(slide, 'Tabla 42').table

    # Año 1 (col 3) is the active header; Nuevos (col 0) turns grey.
    assert t.rows[0].cells[3].fill.fore_color.rgb == main.MATRIX_RED
    assert t.rows[0].cells[0].fill.fore_color.rgb == main.MATRIX_GREY
    # Under Año 1: Semipresencial (offset 1 → col 4) selected, siblings pink.
    assert t.rows[1].cells[4].fill.fore_color.rgb == main.MATRIX_RED
    assert t.rows[1].cells[3].fill.fore_color.rgb == main.MATRIX_PINK
    assert t.rows[1].cells[5].fill.fore_color.rgb == main.MATRIX_PINK
    # The whole Nuevos group stays grey.
    assert t.rows[1].cells[0].fill.fore_color.rgb == main.MATRIX_GREY

    # Unknown modalidad → age matches but modality doesn't; nothing selected.
    assert not main.update_footer_matrix(prs.slides[0], 'Nuevos', 'OTRA')


if __name__ == '__main__':
    test_var_table()
    test_footer_matrix()
    print('ok')
